# Fadil Eldin
# Dec/2024
# Split PDF
# Merge PDF
# Rotate PDF 180 degrees
# WaterMark PDF with png
# JPG to PDF
# PDF to JPG
# encrypt PDF
# unlock PDF
# PDF to Word
# PDF to PPTX
# HTML to PDF
#-----------------------------------------------
def color_picker(color_name):
    color_dict = {
        "Red": (255, 0, 0),
        "Green": (0, 255, 0),
        "Blue": (0, 0, 255),
        "Yellow": (255, 255, 0),
        "Cyan": (0, 255, 255),
        "Magenta": (255, 0, 255),
        "White": (255, 255, 255),
        "Black": (0, 0, 0),
        "Gray": (128, 128, 128),
        "Orange": (255, 165, 0),
        "Purple": (128, 0, 128),
        "Pink": (255, 192, 203),
        "Brown": (165, 42, 42),
        "Turquoise": (64, 224, 208),
        "Gold": (255, 215, 0)
    }
    color_name = color_name.capitalize()
    if color_name in color_dict:
        return color_dict[color_name]
    else:
        return "Color not found in the dictionary."
#-----------------------------------------------
def split_pdf(input_pdf, output_folder):
    #pip install PyPDF2
    import PyPDF2
    import os
    # Create the output folder if it doesn't exist
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    pdf_reader = PyPDF2.PdfReader(open(input_pdf, 'rb'))
    for page_num in range(len(pdf_reader.pages)):
        #Initialize the writer at each page, so that pages are not merged back
        pdf_writer = PyPDF2.PdfWriter()
        pdf_writer.add_page(pdf_reader.pages[page_num])
        output_filename = f'{output_folder}/page_{page_num + 1}.pdf'
        with open(output_filename, 'wb') as output_file:
            pdf_writer.write(output_file)
    print(f"PDF split into {len(pdf_reader.pages)} pages")
#------------------------------------------------------------------
def merge_pdfs(input_pdfs, output_pdf):
    #pip install PyPDF2
    import PyPDF2
    import os
    # Create the output folder if it doesn't exist
    output_folder=os.path.dirname(output_pdf)
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    pdf_merger = PyPDF2.PdfMerger()
    for pdf in input_pdfs:
        pdf_merger.append(pdf)
    with open(output_pdf, 'wb') as output_file:
        pdf_merger.write(output_file)
    print(f"Merged {len(input_pdfs)} PDFs into {output_pdf}")
#------------------------------------------------------------------
def rotate_pdf_180(input_pdf, output_pdf):
    #pip install PyPDF2
    from PyPDF2 import PdfReader, PdfWriter

    import os
    # Create the output folder if it doesn't exist
    output_folder = os.path.dirname(output_pdf)
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Open the input PDF file
    with open(input_pdf, 'rb') as file:
        reader = PdfReader(file)
        writer = PdfWriter()

        # Rotate each page by 180 degrees
        for page in reader.pages:
            page.rotate(180)
            writer.add_page(page)

        # Write the rotated pages to a new PDF file
        with open(output_pdf, 'wb') as output_file:
            writer.write(output_file)
#------------------------------------------------------------------
def add_png_watermark(input_pdf, output_pdf, watermark_png, opacity=1, size_percent=1):
    # pip install PyPDF4 Pillow
    from PyPDF4 import PdfFileReader, PdfFileWriter
    from PIL import Image
    import io
    import os
    # Create the output folder if it doesn't exist
    output_folder = os.path.dirname(output_pdf)
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    pdf_reader = PdfFileReader(open(input_pdf, 'rb'))
    pdf_writer = PdfFileWriter()

    watermark = Image.open(watermark_png).convert("RGBA")

    for page_num in range(pdf_reader.numPages):
        page = pdf_reader.getPage(page_num)
        page_width = page.mediaBox.getWidth()
        page_height = page.mediaBox.getHeight()

        # Resize watermark
        watermark_width = int(page_width * size_percent)
        watermark_height = int(watermark_width * watermark.size[1] / watermark.size[0])
        watermark_resized = watermark.resize((watermark_width, watermark_height), Image.LANCZOS)

        # Apply opacity
        watermark_with_opacity = Image.new("RGBA", watermark_resized.size, (255, 255, 255, 0))
        watermark_with_opacity = Image.blend(watermark_with_opacity, watermark_resized, opacity)

        # Create PDF page with watermark
        watermark_page = io.BytesIO()
        watermark_with_opacity.save(watermark_page, 'PDF')
        watermark_pdf = PdfFileReader(watermark_page)

        # Calculate center position
        x = (page_width - watermark_width) / 2
        y = (page_height - watermark_height) / 2

        # Create a new blank page and merge original content and watermark
        new_page = pdf_writer.addBlankPage(page_width, page_height)
        new_page.mergePage(page)
        new_page.mergeTranslatedPage(watermark_pdf.getPage(0), x, y)

    with open(output_pdf, 'wb') as out:
        pdf_writer.write(out)
    print(f"Watermarked PDF saved as {output_pdf}")
#------------------------------------------------------------------
def add_text_watermark(input_pdf, output_pdf, watermark_text, opacity=0.5, angle=45, color=(0, 0, 0), font_size=50, font_path=None, is_rtl=False):
    #pip install arabic-reshaper python-bidi
    #pip install --upgrade --force-reinstall reportlab
    from PyPDF4 import PdfFileReader, PdfFileWriter
    from reportlab.pdfgen import canvas
    from reportlab.lib.colors import Color
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from bidi.algorithm import get_display
    import arabic_reshaper
    import io
    import os
    # Create the output folder if it doesn't exist
    output_folder = os.path.dirname(output_pdf)
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    reader = PdfFileReader(open(input_pdf, 'rb'))
    writer = PdfFileWriter()

    # Get the size of the first page
    first_page = reader.getPage(0)
    page_width = first_page.mediaBox.getWidth()
    page_height = first_page.mediaBox.getHeight()

    # Register the font if a path is provided
    if font_path:
        font_name = "CustomFont"
        pdfmetrics.registerFont(TTFont(font_name, font_path))
    else:
        font_name = "Helvetica"

    # Handle RTL text if specified
    if is_rtl:
        reshaped_text = arabic_reshaper.reshape(watermark_text)
        watermark_text = get_display(reshaped_text)

    # Create watermark
    packet = io.BytesIO()
    can = canvas.Canvas(packet, pagesize=(page_width, page_height))
    can.setFont(font_name, font_size)
    can.setFillColor(Color(color[0], color[1], color[2], alpha=opacity))
    can.saveState()
    can.translate(page_width/2, page_height/2)
    can.rotate(angle)
    can.drawCentredString(0, 0, watermark_text)
    can.restoreState()
    can.save()
    packet.seek(0)
    watermark = PdfFileReader(packet)

    for i in range(reader.getNumPages()):
        page = reader.getPage(i)
        page.mergePage(watermark.getPage(0))
        writer.addPage(page)

    with open(output_pdf, "wb") as out_file:
        writer.write(out_file)

    print(f"Watermarked PDF saved as {output_pdf}")
#------------------------------------------------------------------
def jpg_to_pdf(jpg_path, pdf_path):
    # pip install Pillow PyPDF2
    from PIL import Image
    import PyPDF2
    import io
    import os
    # Create the output folder if it doesn't exist
    output_folder = os.path.dirname(pdf_path)
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Open the JPG image
    image = Image.open(jpg_path)

    # Convert the image to RGB mode if it's not already
    if image.mode != 'RGB':
        image = image.convert('RGB')

    # Create a PDF writer object
    pdf_writer = PyPDF2.PdfWriter()

    # Convert the image to PDF
    pdf_bytes = io.BytesIO()
    image.save(pdf_bytes, format='PDF')
    pdf_bytes.seek(0)

    # Add the PDF page to the writer
    pdf_reader = PyPDF2.PdfReader(pdf_bytes)
    pdf_writer.add_page(pdf_reader.pages[0])

    # Save the PDF
    with open(pdf_path, 'wb') as output_file:
        pdf_writer.write(output_file)
#------------------------------------------------------------------
def pdf_to_jpg(pdf_path, output_folder, dpi=300):
    # pip install pdf2image
    # Need to install poppler
    # https://github.com/Belval/pdf2image
    # https://github.com/oschwartz10612/poppler-windows
    # add to PATH: C:\Program Files\poppler-24.08.0\Library\bin
    from pdf2image import convert_from_path
    import os
    # Create the output folder if it doesn't exist
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Create the output folder if it doesn't exist
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Convert PDF to list of images
    images = convert_from_path(pdf_path, dpi=dpi)

    # Save each image
    for i, image in enumerate(images):
        image_path = os.path.join(output_folder, f'page_{i+1}.jpg')
        image.save(image_path, 'JPEG')
        print(f"Saved: {image_path}")

    print(f"Converted {len(images)} pages to JPG")
#------------------------------------------------------------------
def encrypt_pdf(input_path, output_path, password):
    # pip install PyPDF2
    from PyPDF2 import PdfReader, PdfWriter
    import os
    # Create the output folder if it doesn't exist
    output_folder = os.path.dirname(output_path)
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Create a PDF reader object
    reader = PdfReader(input_path)

    # Create a PDF writer object
    writer = PdfWriter()

    # Add all pages to the writer
    for page in reader.pages:
        writer.add_page(page)

    # Encrypt the new PDF with the provided password
    writer.encrypt(password)

    # Write the encrypted PDF to a file
    with open(output_path, "wb") as output_file:
        writer.write(output_file)

    print(f"Created encrypted PDF: {output_path}")
# ------------------------------------------------------------------
def unlock_pdf(input_pdf, output_pdf, password):
    # pip install pikepdf
    # Another method in PyPDF2-Unlock.py
    import pikepdf
    import os
    # Create the output folder if it doesn't exist
    output_folder = os.path.dirname(output_pdf)
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    try:
        # Open the protected PDF with the password
        with pikepdf.open(input_pdf, password=password) as pdf:
            # Save the unlocked PDF
            pdf.save(output_pdf)
        print(f"Unlocked PDF saved as: {output_pdf}")
    except pikepdf._qpdf.PasswordError:
        print("Incorrect password.")
    except Exception as e:
        print(f"Failed to unlock PDF: {e}")
# ------------------------------------------------------------------
def pdf_to_word(input_pdf, output_docx):
    # pip install pdf2docx --user
    from pdf2docx import Converter
    c = Converter(input_pdf)
    c.convert(output_docx)
    c.close()
    return
    # pip install aspose-words
    # Aspose.Words is a paid library, and its advanced features may require a license for commercial use.
    # Will add to the PDF :
    # Created with an evaluation copy of Aspose.Words.
    # To remove all limitations, you can use Free Temporary License https://products.aspose.com/words/temporary-license/

    import aspose.words as aw
    import os
    # Create the output folder if it doesn't exist
    output_folder = os.path.dirname(output_docx)
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    doc = aw.Document(input_pdf)
    doc.save(output_docx)
    print(f"Word saved to {output_docx}")
#------------------------------------------------------------------
def pdf_to_pptx(pdf_path, pptx_path):
    # pip install python-pptx pymupdf pillow
    import fitz  # PyMuPDF
    from pptx import Presentation
    from pptx.util import Inches
    from io import BytesIO
    from PIL import Image
    import os
    # Create the output folder if it doesn't exist
    output_folder = os.path.dirname(pptx_path)
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Open the PDF file
    pdf_document = fitz.open(pdf_path)
    presentation = Presentation()

    for page_num in range(len(pdf_document)):
        page = pdf_document[page_num]

        # Convert page to an image
        pix = page.get_pixmap()
        img_data = BytesIO(pix.tobytes("png"))
        img = Image.open(img_data)

        # Create a slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # Blank layout

        # Add image to slide
        left = Inches(0)
        top = Inches(0)
        slide_width = Inches(10)  # Default slide width in PowerPoint
        slide_height = Inches(7.5)  # Default slide height in PowerPoint

        slide.shapes.add_picture(img_data, left, top, slide_width, slide_height)

    # Save the presentation
    presentation.save(pptx_path)
    print(f"Presentation saved to {pptx_path}")
#------------------------------------------------------------------
def html_to_pdf_1(html_content, pdf_path):
    # Method 1: pdfkit
    # pip install pdfkit
    # Windows Download from the official wkhtmltopdf website and add it to your system PATH.
    # Ubuntu: sudo apt-get install wkhtmltopdf
    import pdfkit
    import os
    # Create the output folder if it doesn't exist
    output_folder = os.path.dirname(pdf_path)
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    pdfkit.from_string(html_content, pdf_path)
#--------------------------------------------------
async def html_to_pdf_2(html_content, pdf_path):
    # Method 2: Pyppeteer
    # pip install pyppeteer
    import asyncio
    from pyppeteer import launch
    import os
    # Create the output folder if it doesn't exist
    output_folder = os.path.dirname(pdf_path)
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    browser = await launch()
    page = await browser.newPage()
    await page.setContent(html_content)
    await page.pdf({'path': pdf_path, 'format': 'A4'})
    await browser.close()
#--------------------------------------------------
def html_to_pdf_3(html_content, pdf_path):
    # Method 3: xhtml2pdf
    # pip install xhtml2pdf
    from xhtml2pdf import pisa
    import os
    # Create the output folder if it doesn't exist
    output_folder = os.path.dirname(pdf_path)
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    with open(pdf_path, "wb") as pdf_file:
        pisa_status = pisa.CreatePDF(html_content, dest=pdf_file)
        return not pisa_status.err
#------------------------------------------------------------------
if __name__ == '__main__':
    print('PyCharm')

    # Split
    # input_pdf = 'PDF_TEST/sample.pdf'
    # output_folder = 'PDF_TEST/OUT/SPLIT/'
    # split_pdf(input_pdf, output_folder)
    # exit()
    #==================================================================================================================================
    #  Merge
    # f1='C:/Users/Fadil/Downloads/TEST/Name_Change_Maya_Ryan_1.pdf'
    # f2='C:/Users/Fadil/Downloads/TEST/Name_Change_Maya_Ryan_2.pdf'
    # f3='C:/Users/Fadil/Downloads/TEST/Name_Change_Maya_Ryan_3.pdf'
    #
    # input_pdfs = [f1,f2, f3]
    # output_pdf = 'C:/Users/Fadil/Downloads/TEST/Name_Change.pdf'
    # merge_pdfs(input_pdfs, output_pdf)
    # exit()
    #==================================================================================================================================
    #  Rotate
    # input_pdf = 'PDF_TEST/sample.pdf'
    # output_pdf = 'PDF_TEST/OUT/rotated-sample.pdf'
    # rotate_pdf_180(input_pdf, output_pdf)
    # exit()
    # ==================================================================================================================================
    #  PNG Watermark
    # input_pdf = 'PDF_TEST/sample.pdf'
    # watermark_png = 'PDF_TEST/confidential.png'
    # watermarked_pdf = 'PDF_TEST/OUT/png_watermarked_pdf.pdf'
    # add_png_watermark(input_pdf, watermarked_pdf, watermark_png,opacity=0.32, size_percent=0.8)
    # exit()
    #==================================================================================================================================
    # Do not uncomment next 2 lines
    # input_color='Red'
    # rgb_color = color_picker(input_color)

    # Text Watermark
    # input_pdf = 'PDF_TEST/sample.pdf'
    # watermarked_pdf = 'PDF_TEST/OUT/text_watermarked_pdf.pdf'
    # add_text_watermark(input_pdf, watermarked_pdf, "Confidential",
    #                    opacity = 0.33,
    #                    angle = 45,
    #                    color = rgb_color,
    #                    font_size=130  # Font size in points
    #                    )
    # exit()
    #-----------------------------------------------------------------
    # input_pdf = 'PDF_TEST/sample.pdf'
    # watermarked_pdf = 'PDF_TEST/OUT/arabic_watermarked_pdf.pdf'
    # arabic_font_path ='C:/Windows/Fonts/trado.ttf'
    # add_text_watermark(
    #     input_pdf,
    #     watermarked_pdf,
    #     "سري للغاية",  # Arabic text for "Top Secret"
    #     opacity=0.4,
    #     angle=45,
    #     color = rgb_color,
    #     font_size=200,
    #     font_path=arabic_font_path,
    #     is_rtl=True  # Specify that this is RTL text
    # )
    # exit()
    # ==================================================================================================================================
    # JPG to PDF
    # jpg_path = 'C:/Users/Fadil/Downloads/TEST/Name_Change_Maya_Ryan_1.jpg'
    # pdf_path = 'C:/Users/Fadil/Downloads/TEST/Name_Change_Maya_Ryan_1.pdf'
    # jpg_to_pdf(jpg_path, pdf_path)
    # exit()
    #==================================================================================================================================
    # PDF to JPG
    pdf_path = 'C:/Users/Fadil/Downloads/vegas.pdf'
    output_folder = 'C:/Users/Fadil/Downloads/'
    pdf_to_jpg(pdf_path, output_folder)
    exit()
    #==================================================================================================================================
    # Password Encrypt
    # input_pdf = 'PDF_TEST/sample.pdf'
    # output_pdf = 'PDF_TEST/OUT/encrypted_test.pdf'
    # password = 'mypass'
    # encrypt_pdf(input_pdf, output_pdf, password)
    # exit()
    # ==================================================================================================================================
    # Unlock PDF
    # encrypted_pdf = 'PDF_TEST/OUT/encrypted_test.pdf'
    # unlocked_pdf = 'PDF_TEST/OUT/unlocked_test.pdf'
    # password = 'mypass'
    # unlock_pdf(encrypted_pdf, unlocked_pdf, password)
    # exit()
    #==================================================================================================================================
    #  PDF To Word
    pdf_path = 'C:/Users/Fadil/Downloads/Business Unit Leader Job Description.pdf'
    word_path = 'C:/Users/Fadil/Downloads/Business Unit Leader Job Description.docx'
    pdf_to_word(pdf_path, word_path)
    exit()
    # ==================================================================================================================================
    #  PDF To PPTX
    # pdf_path = 'PDF_TEST/sample.pdf'
    # pptx_path = 'PDF_TEST/OUT/sample.pptx'
    # pdf_to_pptx(pdf_path, pptx_path)
    # exit()
# pip freeze > requirements.txt